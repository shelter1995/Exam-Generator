"""
多文件类型解析器
统一入口：parse(file_path) -> {texts, metadatas, ids}
"""
import uuid
import logging
from pathlib import Path
from typing import Dict, List, Any

import config

logger = logging.getLogger(__name__)


def chunk_text(text: str) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + config.CHUNK_SIZE
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - config.CHUNK_OVERLAP
    return chunks


def parse(file_path: str, db_id: str = "default") -> Dict[str, Any]:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in config.SUPPORTED_DOCUMENTS:
        return _parse_document(path, db_id)
    elif ext in config.SUPPORTED_AUDIO:
        return _parse_audio(path, db_id)
    elif ext in config.SUPPORTED_VIDEO:
        return _parse_video(path, db_id)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _parse_document(path: Path, db_id: str) -> Dict[str, Any]:
    ext = path.suffix.lower()
    full_text = ""

    if ext == ".pdf":
        full_text = _extract_pdf(path)
    elif ext == ".docx":
        full_text = _extract_docx(path)
    elif ext == ".xlsx":
        full_text = _extract_xlsx(path)
    elif ext == ".pptx":
        full_text = _extract_pptx(path)
    elif ext in (".txt", ".md", ".py", ".js", ".html", ".json"):
        full_text = _extract_text(path)

    return _chunk_and_meta(full_text, path.name, ext, db_id)


def _extract_pdf(path: Path) -> str:
    try:
        import pypdf
        texts = []
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception:
        try:
            import pdfplumber
            texts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
            return "\n".join(texts)
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            return ""


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        logger.error(f"DOCX解析失败: {e}")
        return ""


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path)
        texts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                texts.append(f"【Sheet: {sheet_name}】\n" + "\n".join(rows))
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"XLSX解析失败: {e}")
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"【第{i+1}页】\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PPTX解析失败: {e}")
        return ""


def _extract_text(path: Path) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    logger.error(f"文本解码失败: {path}")
    return ""


def _parse_audio(path: Path, db_id: str) -> Dict[str, Any]:
    """音频：Whisper 转录"""
    try:
        import whisper
        logger.info(f"开始转录音频: {path.name}")
        # 优先使用本地模型
        model_path = config.WHISPER_MODEL_PATH or "base"
        model = whisper.load_model(model_path)
        result = model.transcribe(str(path), language="zh", fp16=False)
        text = result["text"]
        return _chunk_and_meta(text, path.name, path.suffix, db_id)
    except Exception as e:
        logger.error(f"音频解析失败: {e}")
        return {"texts": [], "metadatas": [], "ids": []}


def _parse_video(path: Path, db_id: str) -> Dict[str, Any]:
    """视频：提取音频 -> Whisper 转录"""
    import tempfile
    try:
        import subprocess
        logger.info(f"开始处理视频: {path.name}")

        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            audio_path = tmp.name

        ffmpeg_cmd = config.FFMPEG_PATH or "ffmpeg"
        cmd = [
            ffmpeg_cmd, "-i", str(path), "-vn", "-acodec", "pcm_s16le",
            "-ar", "16000", "-ac", "1", "-y", audio_path
        ]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)

        import whisper
        model_path = config.WHISPER_MODEL_PATH or "base"
        model = whisper.load_model(model_path)
        result = model.transcribe(audio_path, language="zh", fp16=False)
        text = result["text"]

        Path(audio_path).unlink(missing_ok=True)
        return _chunk_and_meta(text, path.name, path.suffix, db_id)
    except Exception as e:
        logger.error(f"视频解析失败: {e}")
        return {"texts": [], "metadatas": [], "ids": []}


def _chunk_and_meta(text: str, filename: str, ext: str, db_id: str) -> Dict[str, Any]:
    chunks = chunk_text(text) if text else []
    texts = []
    metadatas = []
    ids = []
    for i, chunk in enumerate(chunks):
        texts.append(chunk)
        metadatas.append({
            "source": filename,
            "type": ext.lstrip("."),
            "chunk": i + 1,
            "db_id": db_id
        })
        ids.append(str(uuid.uuid4()))
    return {"texts": texts, "metadatas": metadatas, "ids": ids}
